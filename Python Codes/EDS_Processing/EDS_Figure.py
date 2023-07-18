from PIL import Image, ImageDraw, ImageFilter
import os
import numpy as np
import matplotlib.pyplot as plt
import matplotlib as mpl
import pandas as pd
import glob
import scipy.optimize as optimize
import pylab
from pptx import Presentation

#report_dir = r'C:\Users\benja\OneDrive - Northeastern University\Gallaway Group\SEM EDS\2023\0711\MDB_Flake_-_Cycled\Report Exports'
slash = '\\'
print(os.getcwd())
#os.chdir(report_dir)
print(os.getcwd())
print(os.listdir(os.getcwd()))

#report_file_names = os.listdir(report_dir)

report_1 = Presentation('MDB_Flake_Cycled_G_Area_1_Report.pptx')

print('presentation loaded')

for slidecount, slide in enumerate(report_1.slides):
    print('slide %s' % slidecount)
    for shape in slide.shapes:
        if shape.shape_type == 17:
            print('%s : %s' % (shape.shape_type, shape.text))
        else:
            print('%s : %s' % (shape.shape_type, shape.name))


print('exploring slide 1')
if __name__ == '__main__':
    print('main script')