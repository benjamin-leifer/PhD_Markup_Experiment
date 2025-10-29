"""
Recreate the combined anode + cathode microstructure figure as a 1-slide PowerPoint.

- Anode (left, Cu collector): 9x4 particles, each fully filled with graphite “flakes”
  (two randomized chord directions per particle), plus a PVDF squiggle on ~1/3 perimeter.
- Cathode (right, Al collector): 9x4 particles, each host disk tiled with tiny circles
  (primary particles), rim carbon black (~30% perimeter with ~20% D dots), plus PVDF squiggle.
- Leftmost anode column is tangent to Cu; rightmost cathode column is tangent to Al.
- Center separator. All geometry stable; randomness is seeded for reproducibility.

Tip for animation: shapes are named like:
  "anode_r{r}_c{c}_base", "anode_r{r}_c{c}_flake_*", "anode_r{r}_c{c}_pvdf",
  "cathode_r{r}_c{c}_host", "cathode_r{r}_c{c}_prim_*", "cathode_r{r}_c{c}_cb_*"
so you can target them with PowerPoint animations.
"""

from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import math, random

# ----------------------- CONFIG (edit these) -----------------------
OUT_PATH = r"C:\Users\benja\Downloads\10_22_2025 Temp\Slide Trials\Local_Images\Cell_Combined_AnodeCathode_9x4.pptx"
SEED = 20251027            # change per frame to re-randomize (for animations)
ROWS, COLS = 9, 4
ROW_GAP = 0.06             # vertical gap between rows (inches)
COL_GAP = 0.08             # horizontal gap between columns (inches)

# PVDF arc parameters (both sides)
PVDF_FRACTION = 1/3        # ~one-third perimeter
PVDF_WIDTH = 0.045         # inches
PVDF_WIGGLE_AMP_FRACTION = 0.05  # radial wiggle amplitude as fraction of particle radius
PVDF_WIGGLE_FREQ = 4.0

# Graphite flake parameters (anode)
FLAKE_WIDTH_RATIO = 0.12   # thickness of a “flake linelet” relative to particle diameter
# Cathode internal primary particles
CATHODE_PRIM_COUNT = 12
CATHODE_PRIM_SCALE = (0.08, 0.16)  # primary radius range relative to big-particle radius
# Rim carbon black (cathode)
CB_COVERAGE = 0.30         # ~30% of rim covered
CB_SIZE_REL = 0.20         # each CB dot ≈ 20% of particle diameter

# Canvas geometry
FRAME = (0.5, 0.7, 9.0, 6.0)  # x, y, w, h (inches)
CC_W = 0.30
SEP_W = 0.50
MARGIN = 0.10               # inner padding of frame
SEP_GAP = 0.05              # keep grid a hair away from separator
# ------------------------------------------------------------------

# Colors
FLAKE_DARK  = RGBColor(105, 115, 130)
FLAKE_LIGHT = RGBColor(165, 175, 190)
EDGE        = RGBColor(120, 130, 145)
PVDF_COL    = RGBColor(235, 210, 125)

PURPLE_HOST = RGBColor(165, 110, 205)
PURPLE_CORE = RGBColor(150, 95, 190)
PURPLE_DOT  = RGBColor(185, 135, 220)
CBLACK      = RGBColor(30, 30, 30)

AL          = RGBColor(170, 170, 170)
CU          = RGBColor(186, 120, 60)
SEPBK       = RGBColor(238, 238, 238)
FRAME_COL   = RGBColor(200, 200, 200)

rnd = random.Random(SEED)

def add_rect(slide, name, x, y, w, h, fill=None, line=None, rot=0, rounded=False):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    shp.name = name
    if rounded:
        try: shp.adjustments[0] = 0.22
        except: pass
    if fill is None: shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line
    if rot: shp.rotation = rot
    return shp

def add_oval(slide, name, x, y, w, h, fill, line=None):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.name = name
    shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line
    return shp

def pvdf_wavy_arc(slide, name_prefix, cx, cy, R, sweep_fraction=1/3, width=0.045,
                  amp_frac=0.05, freq=4.0, start_deg=None):
    """One curved 'squiggle' along a rim arc."""
    if start_deg is None:
        start_deg = rnd.uniform(0, 360.0 - 360.0*sweep_fraction)
    sweep_deg = 360.0 * sweep_fraction
    steps = 14
    for i in range(steps-1):
        t  = math.radians(start_deg + sweep_deg*(i/(steps-1)))
        tn = math.radians(start_deg + sweep_deg*((i+1)/(steps-1)))
        rr  = R * (1.02 + amp_frac * math.sin(2*math.pi*freq*(i/(steps-1))))
        rrn = R * (1.02 + amp_frac * math.sin(2*math.pi*freq*((i+1)/(steps-1))))
        x  = cx + rr  * math.cos(t)
        y  = cy + rr  * math.sin(t)
        xn = cx + rrn * math.cos(tn)
        yn = cy + rrn * math.sin(tn)
        L = ((xn-x)**2 + (yn-y)**2)**0.5
        if L <= 1e-6:
            continue
        theta = math.degrees(math.atan2(yn-y, xn-x))
        add_rect(slide, f"{name_prefix}_pvdf_seg_{i}", x, y - width/2.0, L, width,
                 fill=PVDF_COL, line=None, rot=theta, rounded=True)

def flake_texture(slide, name_prefix, cx, cy, D, width_ratio=0.12):
    """Two heavy chord passes at randomized angles; full interior coverage."""
    R = D/2.0
    base = rnd.uniform(-50, 50)
    angles = (base, base + rnd.uniform(60, 110))
    add_oval(slide, f"{name_prefix}_base", cx - R, cy - R, 2*R, 2*R, FLAKE_DARK, EDGE)
    w = D * width_ratio
    cols = [FLAKE_LIGHT, FLAKE_DARK]
    for pass_idx, (ang, col) in enumerate(zip(angles, cols)):
        theta = math.radians(ang)
        vx, vy = -math.sin(theta), math.cos(theta)
        step = w * 0.9
        t = -R
        k = 0
        while t <= R:
            if abs(t) <= R:
                L = 2.0 * math.sqrt(max(0.0, R*R - t*t))
                px = cx + vx * t
                py = cy + vy * t
                add_rect(slide, f"{name_prefix}_flake_{pass_idx}_{k}",
                         px - L/2.0, py - w/2.0, L, w, fill=col, line=None, rot=ang, rounded=True)
                k += 1
            t += step

def cathode_particle(slide, name_prefix, cx, cy, D,
                     prim_count=12, prim_scale=(0.08, 0.16),
                     cb_coverage=0.30, cb_size_rel=0.20):
    R = D/2.0
    add_oval(slide, f"{name_prefix}_host", cx - R, cy - R, 2*R, 2*R, PURPLE_HOST, None)
    add_oval(slide, f"{name_prefix}_core", cx - R*0.92, cy - R*0.92, 2*R*0.92, 2*R*0.92, PURPLE_CORE, None)
    # primary dots
    for i in range(prim_count):
        rp = rnd.uniform(prim_scale[0]*R, prim_scale[1]*R)
        placed = False
        for _try in range(30):
            dx = rnd.uniform(-R*0.85, R*0.85)
            dy = rnd.uniform(-R*0.85, R*0.85)
            if (dx*dx + dy*dy) <= (R - rp*1.05)**2:
                add_oval(slide, f"{name_prefix}_prim_{i}", cx + dx - rp, cy + dy - rp,
                         2*rp, 2*rp, PURPLE_DOT, None)
                placed = True
                break
        if not placed:
            # fallback near center
            add_oval(slide, f"{name_prefix}_prim_{i}", cx - rp, cy - rp, 2*rp, 2*rp, PURPLE_DOT, None)
    # rim CB dots
    d_cb = cb_size_rel * D
    N = max(2, int(round((cb_coverage * 2*math.pi*R) / d_cb)))
    used = []
    min_sep = max(0.8 * d_cb / R, 0.25)
    tries = 0
    while len(used) < N and tries < 250:
        th = rnd.uniform(0, 2*math.pi)
        if all(abs((th - t + math.pi)%(2*math.pi)-math.pi) > min_sep for t in used):
            used.append(th)
        tries += 1
    for i, th in enumerate(used):
        rr = R*0.96
        x = cx + rr*math.cos(th); y = cy + rr*math.sin(th)
        add_oval(slide, f"{name_prefix}_cb_{i}", x - d_cb/2.0, y - d_cb/2.0, d_cb, d_cb, CBLACK, None)

def main():
    prs = Presentation()
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Frame & scaffolding
    fx, fy, fw, fh = FRAME
    add_rect(s, "frame", fx, fy, fw, fh, None, FRAME_COL, rounded=True)
    ccx_L = fx + MARGIN
    ccx_R = fx + fw - MARGIN - CC_W
    ccy   = fy + MARGIN
    cch   = fh - 2*MARGIN

    add_rect(s, "collector_Cu", ccx_L, ccy, CC_W, cch, CU, None)
    add_rect(s, "collector_Al", ccx_R, ccy, CC_W, cch, AL, None)

    sep_x = fx + (fw - SEP_W)/2.0
    add_rect(s, "separator", sep_x, ccy, SEP_W, cch, SEPBK, None)

    # Common particle size from vertical packing
    Hgrid = cch
    D = (Hgrid - (ROWS - 1)*ROW_GAP) / ROWS
    R = D/2.0
    row_centers = [ccy + R + r*(D + ROW_GAP) for r in range(ROWS)]

    # --- Anode grid (left of separator), tangent to Cu ---
    gridL_x1 = ccx_L + CC_W
    an_cols = [gridL_x1 + R + c*(D + COL_GAP) for c in range(COLS)]
    max_an_right = an_cols[-1] + R
    if max_an_right > sep_x - SEP_GAP:
        span = (sep_x - SEP_GAP) - (an_cols[0] - R)
        D = span / COLS
        R = D/2.0
        an_cols = [gridL_x1 + R + c*D for c in range(COLS)]
        row_centers = [ccy + R + r*(D + ROW_GAP) for r in range(ROWS)]

    for r, yy in enumerate(row_centers):
        for c, xx in enumerate(an_cols):
            name_prefix = f"anode_r{r}_c{c}"
            flake_texture(s, name_prefix, xx, yy, D, width_ratio=FLAKE_WIDTH_RATIO)
            pvdf_wavy_arc(
                s, name_prefix, xx, yy, D/2.0,
                sweep_fraction=PVDF_FRACTION,
                width=PVDF_WIDTH,
                amp_frac=PVDF_WIGGLE_AMP_FRACTION,
                freq=PVDF_WIGGLE_FREQ,
                start_deg=None  # randomized
            )

    # --- Cathode grid (right of separator), tangent to Al ---
    gridR_x2 = ccx_R
    ca_cols = [gridR_x2 - R - c*(D + COL_GAP) for c in range(COLS)]
    min_ca_left = ca_cols[-1] - R
    if min_ca_left < sep_x + SEP_W + SEP_GAP:
        span = (ca_cols[0] + R) - (sep_x + SEP_W + SEP_GAP)
        D = span / COLS
        R = D/2.0
        ca_cols = [gridR_x2 - R - c*D for c in range(COLS)]
        row_centers = [ccy + R + r*(D + ROW_GAP) for r in range(ROWS)]

    for r, yy in enumerate(row_centers):
        for c, xx in enumerate(ca_cols):
            name_prefix = f"cathode_r{r}_c{c}"
            cathode_particle(
                s, name_prefix, xx, yy, D,
                prim_count=CATHODE_PRIM_COUNT,
                prim_scale=CATHODE_PRIM_SCALE,
                cb_coverage=CB_COVERAGE,
                cb_size_rel=CB_SIZE_REL
            )
            pvdf_wavy_arc(
                s, name_prefix, xx, yy, D/2.0,
                sweep_fraction=PVDF_FRACTION,
                width=PVDF_WIDTH,
                amp_frac=PVDF_WIGGLE_AMP_FRACTION,
                freq=PVDF_WIGGLE_FREQ,
                start_deg=None  # randomized
            )

    prs.save(OUT_PATH)
    print(f"Saved: {OUT_PATH}")

if __name__ == "__main__":
    main()


"""
Prompt:
“Regenerate it fast” prompt (keep this handy)

Generate a single-slide PowerPoint (.pptx) with a combined anode (left) and cathode (right) microstructure:

Canvas: frame (0.5, 0.7, 9.0, 6.0) in, Cu collector on left (w=0.30 in), Al collector on right (w=0.30 in), separator centered (w=0.50 in). Keep ~0.10 in margins; keep the particle grids ~0.05 in away from the separator.

Anode grid (left of separator): 9×4 circles. The leftmost column is tangent to the Cu collector. Determine particle diameter D from vertical packing with ROW_GAP = 0.06 in and COL_GAP = 0.08 in. Each particle interior is fully filled with graphite “flakes”: draw two overlapping chord passes at randomized angles per particle (linelet thickness = 0.12·D; colors #697382 / #A5AFBE). Add one PVDF squiggle hugging ~⅓ of circumference, thickness 0.045 in, curved with gentle radial wiggle (amp = 0.05·R, freq ≈ 4).

Cathode grid (right of separator): 9×4 circles. The rightmost column is tangent to the Al collector. Inside each particle, tile many small circles (“primary particles”) with radius in [0.08·R, 0.16·R]. Add rim carbon black (coverage ~30% of perimeter) with dot diameter ≈ 0.20·D (randomly spaced, non-conformal). Add the same PVDF squiggle spec as the anode (random start angle per particle).

Colors: Cu = #BA783C, Al = #AAAAAA, separator = #EEEEEE, anode flakes = #697382/#A5AFBE, cathode host = #A56ECD/#965FBE with primaries #B987DC, CB = #1E1E1E, PVDF = #EBD27D.

Deterministic randomness: set seed = 20251027 so I can re-run and produce new frames by changing the seed.

Name shapes so I can animate them later:

anode_r{r}_c{c}_base, anode_r{r}_c{c}_flake_*, anode_r{r}_c{c}_pvdf_*

cathode_r{r}_c{c}_host, cathode_r{r}_c{c}_core, cathode_r{r}_c{c}_prim_*, cathode_r{r}_c{c}_cb_*, cathode_r{r}_c{c}_pvdf_*

Export to Cell_Combined_AnodeCathode_9x4.pptx.

Quick animation workflow (optional, but handy)

Duplicate the slide for each “beat” (e.g., add PVDF, then CB, then separator, etc.).

In Selection Pane, target named shapes (e.g., anode_r3_c1_pvdf_*) and apply Appear/Wipe.

For particle-by-particle reveals, step through rows or columns in order by using the name pattern.

Want me to add a tiny CLI wrapper that increments the seed and outputs ..._frame_001.pptx, ..._frame_002.pptx etc.?
"""