"""
Generate a one-slide PowerPoint graphical abstract for:
"Enabling THF-based weakly solvating electrolytes for extreme low-temperature lithium-ion batteries"

Requires:
  pip install python-pptx

Run:
  python make_graphical_abstract.py

Output:
  graphical_abstract_THF_WSE.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.enum.dml import MSO_LINE

# ----------------------------
# Helpers
# ----------------------------
def rgb(hexstr: str) -> RGBColor:
    hexstr = hexstr.strip("#")
    return RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))

def add_round_box(slide, x, y, w, h, title, body_lines,
                  header_size=14, body_size=11,
                  border_color=rgb("404040"), fill_color=rgb("FFFFFF"),
                  rounding=True):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounding else MSO_SHAPE.RECTANGLE,
        x, y, w, h
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1.25)

    # Text frame
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.06)

    # Header
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.bold = True
    run.font.size = Pt(header_size)
    run.font.name = "Calibri"
    run.font.color.rgb = rgb("000000")
    p.space_after = Pt(4)
    p.alignment = PP_ALIGN.LEFT

    # Body
    for i, line in enumerate(body_lines):
        p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(body_size)
        p.font.name = "Calibri"
        p.font.color.rgb = rgb("000000")
        p.level = 0
        p.space_after = Pt(2)

    return shape



def add_arrow(slide, x1, y1, x2, y2, line_color=rgb("404040"), width_pt=2.0):
    # Straight connector with arrow
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    conn.line.color.rgb = line_color
    conn.line.width = Pt(width_pt)
    # set an arrowhead (python-pptx accepts boolean assignment in many versions)
    try:
        conn.line.end_arrowhead = True
    except Exception:
        # older/newer pptx versions may require explicit enum; ignore if not supported
        pass
    return conn

def add_icon_snowflake(slide, cx, cy, size):
    # Simple snowflake: 3 crossing lines
    half = size / 2
    # Horizontal
    l1 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx - half, cy, cx + half, cy)
    # Vertical
    l2 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx, cy - half, cx, cy + half)
    # Diagonal
    l3 = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, cx - half, cy - half, cx + half, cy + half)

    for l in (l1, l2, l3):
        l.line.color.rgb = rgb("606060")
        l.line.width = Pt(1.5)

def add_icon_battery(slide, x, y, w, h):
    # Battery body
    body = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    body.fill.solid()
    body.fill.fore_color.rgb = rgb("FFFFFF")
    body.line.color.rgb = rgb("606060")
    body.line.width = Pt(1.25)
    # Battery terminal
    term_w = w * 0.12
    term_h = h * 0.35
    term = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x + w, y + (h - term_h) / 2, term_w, term_h)
    term.fill.solid()
    term.fill.fore_color.rgb = rgb("FFFFFF")
    term.line.color.rgb = rgb("606060")
    term.line.width = Pt(1.25)

def add_icon_flask(slide, x, y, w, h, label="THF"):
    # Use a simple trapezoid-ish look via "CAN" (closest clean shape) + label
    flask = slide.shapes.add_shape(MSO_SHAPE.TRAPEZOID, x, y, w, h)
    flask.fill.solid()
    flask.fill.fore_color.rgb = rgb("FFFFFF")
    flask.line.color.rgb = rgb("606060")
    flask.line.width = Pt(1.25)

    # Label
    tb = slide.shapes.add_textbox(x, y + h + Inches(0.03), w, Inches(0.25))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.name = "Calibri"
    p.font.color.rgb = rgb("404040")
    p.alignment = PP_ALIGN.CENTER

def add_green_red_paths(slide, x, y, w, h):
    """
    Draw two slim arrows inside a box area:
      - green desired path on top
      - red failure path below
    """
    # Coordinates within region
    pad_x = Inches(0.10)
    start_x = x + pad_x
    end_x = x + w - pad_x
    top_y = y + Inches(0.55)
    bot_y = y + Inches(1.20)

    # Green arrow
    a1 = add_arrow(slide, start_x, top_y, end_x, top_y, line_color=rgb("2E7D32"), width_pt=2.25)
    # Red arrow
    a2 = add_arrow(slide, start_x, bot_y, end_x, bot_y, line_color=rgb("C62828"), width_pt=2.25)

    # Labels
    def label(text, lx, ly):
        t = slide.shapes.add_textbox(lx, ly - Inches(0.14), Inches(1.9), Inches(0.35))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(10)
        p.font.name = "Calibri"
        p.font.color.rgb = rgb("000000")
        p.alignment = PP_ALIGN.LEFT

    label("Li⁺ → desolvation → LiC₆", start_x, top_y)
    label("Stable staging", start_x, top_y + Inches(0.18))

    label("Li⁺(THF)ₓ → cointercalation", start_x, bot_y)
    label("Exfoliation / impedance rise / CE loss", start_x, bot_y + Inches(0.18))

def add_small_caption(slide, x, y, w, text):
    tb = slide.shapes.add_textbox(x, y, w, Inches(0.22))
    tf = tb.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(9)
    p.font.name = "Calibri"
    p.font.color.rgb = rgb("404040")
    p.alignment = PP_ALIGN.CENTER

# ----------------------------
# Build presentation
# ----------------------------
prs = Presentation()
# Use a wide-ish canvas: default is 10 x 7.5; we'll use a single slide and layout within it
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# Title
title = slide.shapes.add_textbox(Inches(0.25), Inches(0.15), Inches(9.5), Inches(0.5))
tf = title.text_frame
tf.clear()
p = tf.paragraphs[0]
p.text = "Enabling THF-based weakly solvating electrolytes for extreme low-temperature lithium-ion batteries"
p.font.size = Pt(18)
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.CENTER

# Geometry
y0 = Inches(1.0)
h = Inches(2.4)

x1 = Inches(0.3); w1 = Inches(2.1)
gap = Inches(0.25)
x2 = x1 + w1 + gap; w2 = Inches(2.1)
x3 = x2 + w2 + gap; w3 = Inches(2.45)
x4 = x3 + w3 + gap; w4 = Inches(2.2)

# Boxes
b1 = add_round_box(
    slide, x1, y0, w1, h,
    "NEED",
    [
        "Extreme cold (≤ −40 to −51 °C)",
        "Power/capacity loss in carbonates",
        "High viscosity / phase limits",
        "Interfacial kinetics dominate",
    ],
)

b2 = add_round_box(
    slide, x2, y0, w2, h,
    "OPPORTUNITY",
    [
        "THF-based WSE",
        "Low viscosity → fast transport",
        "Weak solvation → faster Li⁺ desolvation",
    ],
)

b3 = add_round_box(
    slide, x3, y0, w3, h,
    "CHALLENGE: GRAPHITE COMPATIBILITY",
    [
        "Graphite instability limits THF systems",
    ],
    header_size=13,
    body_size=11
)

b4 = add_round_box(
    slide, x4, y0, w4, h,
    "THESIS APPROACH",
    [
        "Diagnose → Stabilize",
        "¹H/⁷Li NMR: solvation / ion pairing",
        "EIS + DRT (−51 °C): rate limits",
        "XRD / microscopy: staging & exfoliation",
        "Additives (FEC/VC): protective SEI",
        "Formation + exchange: suppress cointercalation",
    ],
    header_size=13,
    body_size=10
)

# Arrows between boxes (centered vertically)
mid_y = y0 + h / 2
add_arrow(slide, x1 + w1, mid_y, x2, mid_y)
add_arrow(slide, x2 + w2, mid_y, x3, mid_y)
add_arrow(slide, x3 + w3, mid_y, x4, mid_y)

# Small arrow captions
add_small_caption(slide, x1 + w1 + Inches(0.02), mid_y - Inches(0.55), gap - Inches(0.04), "electrolyte\ndesign")
add_small_caption(slide, x2 + w2 + Inches(0.02), mid_y - Inches(0.55), gap - Inches(0.04), "graphite\nconstraint")
add_small_caption(slide, x3 + w3 + Inches(0.02), mid_y - Inches(0.55), gap - Inches(0.04), "mechanism-\nguided fixes")

# Icons in Block 1 (snowflake + battery)
add_icon_snowflake(slide, x1 + Inches(0.32), y0 + Inches(0.55), Inches(0.35))
add_icon_battery(slide, x1 + Inches(0.60), y0 + Inches(0.43), Inches(0.40), Inches(0.26))

# Icon in Block 2 (flask)
add_icon_flask(slide, x2 + Inches(0.25), y0 + Inches(0.40), Inches(0.35), Inches(0.42), label="THF")

# Inside Block 3: add desired/failure arrows + labels + warning icon
add_green_red_paths(slide, x3, y0, w3, h)

# A simple warning triangle
warn = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, x3 + w3 - Inches(0.50), y0 + Inches(0.38), Inches(0.30), Inches(0.28))
warn.fill.solid()
warn.fill.fore_color.rgb = rgb("FFFFFF")
warn.line.color.rgb = rgb("C62828")
warn.line.width = Pt(1.5)

# Outcome badge (bottom-right under Block 4)
badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x4 + Inches(0.40), y0 + h + Inches(0.15), Inches(1.8), Inches(0.95))
badge.fill.solid()
badge.fill.fore_color.rgb = rgb("FFFFFF")
badge.line.color.rgb = rgb("404040")
badge.line.width = Pt(1.25)

tf = badge.text_frame
tf.clear()
tf.word_wrap = True
tf.margin_left = Inches(0.10)
tf.margin_top = Inches(0.06)
p = tf.paragraphs[0]
p.text = "OUTCOME"
p.font.size = Pt(12)
p.font.bold = True
p.font.name = "Calibri"
p.alignment = PP_ALIGN.LEFT

for line in ["✓ Stable Gr|NMC cycling", "✓ Rate capability", "✓ Lower impedance growth at −51 °C"]:
    q = tf.add_paragraph()
    q.text = line
    q.font.size = Pt(10)
    q.font.name = "Calibri"
    q.alignment = PP_ALIGN.LEFT

# Add small snowflake icon to outcome badge
add_icon_snowflake(slide, x4 + Inches(0.52), y0 + h + Inches(0.40), Inches(0.22))

# Save
out_file = "graphical_abstract_THF_WSE.pptx"
prs.save(out_file)
print(f"Saved: {out_file}")
